from __future__ import annotations

import csv
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from .models import ExtractedText


def extract_text(path: Path, extension: str) -> list[ExtractedText]:
    extension = extension.lower()
    if extension == ".pdf":
        return _extract_pdf(path)
    if extension == ".docx":
        return _extract_docx(path)
    if extension == ".pptx":
        return _extract_pptx(path)
    if extension == ".xlsx":
        return _extract_xlsx(path)
    if extension == ".csv":
        return _extract_csv(path)
    if extension == ".txt":
        return _extract_txt(path)
    return []


def _extract_pdf(path: Path) -> list[ExtractedText]:
    reader = PdfReader(str(path))
    chunks = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            chunks.append(ExtractedText(text=text, location=f"第 {index} 頁"))
    return chunks


def _extract_docx(path: Path) -> list[ExtractedText]:
    doc = Document(str(path))
    chunks = []
    for index, paragraph in enumerate(doc.paragraphs, start=1):
        text = paragraph.text.strip()
        if text:
            chunks.append(ExtractedText(text=text, location=f"段落 {index}"))
    for table_index, table in enumerate(doc.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            text = " ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if text:
                chunks.append(ExtractedText(text=text, location=f"表格 {table_index} 列 {row_index}"))
    return chunks


def _extract_pptx(path: Path) -> list[ExtractedText]:
    prs = Presentation(str(path))
    chunks = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        if texts:
            chunks.append(ExtractedText(text="\n".join(texts), location=f"投影片 {slide_index}"))
    return chunks


def _extract_xlsx(path: Path) -> list[ExtractedText]:
    wb = load_workbook(str(path), read_only=True, data_only=True)
    chunks = []
    for sheet in wb.worksheets:
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = [str(value) for value in row if value is not None and str(value).strip()]
            if values:
                chunks.append(
                    ExtractedText(
                        text=" ".join(values),
                        location=f"工作表 {sheet.title} 第 {row_index} 列",
                    )
                )
    return chunks


def _extract_csv(path: Path) -> list[ExtractedText]:
    chunks = []
    with path.open("r", encoding="utf-8-sig", newline="", errors="replace") as handle:
        reader = csv.reader(handle)
        for row_index, row in enumerate(reader, start=1):
            text = " ".join(cell.strip() for cell in row if cell.strip())
            if text:
                chunks.append(ExtractedText(text=text, location=f"CSV 第 {row_index} 列"))
    return chunks


def _extract_txt(path: Path) -> list[ExtractedText]:
    chunks = []
    with path.open("r", encoding="utf-8-sig", errors="replace") as handle:
        for line_index, line in enumerate(handle, start=1):
            text = line.strip()
            if text:
                chunks.append(ExtractedText(text=text, location=f"第 {line_index} 行"))
    return chunks

